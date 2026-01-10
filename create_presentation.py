"""
Generate HyperCOT-MS Presentation for Advisor Discussion
Revised version with improved grammar and structure
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from pathlib import Path

BASE_PATH = Path(__file__).parent.resolve()


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_formatted_text(paragraph, text, font_size, is_highlight=False):
    """Add text with inline **bold** formatting to a paragraph using runs."""
    import re
    # Split text by **bold** patterns
    parts = re.split(r'(\*\*.*?\*\*)', text)

    for j, part in enumerate(parts):
        if not part:
            continue
        if part.startswith('**') and part.endswith('**'):
            # Bold text
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.size = Pt(font_size)
            if is_highlight:
                run.font.color.rgb = RGBColor(0, 112, 192)
        else:
            # Normal text
            run = paragraph.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            if is_highlight:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 112, 192)


def add_content_slide(prs, title, content_items=None, image_path=None, image_pos=None,
                      text_width=None, font_size=18):
    """Add a content slide with title, bullet points, and optional image."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Add underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.85), Inches(9.4), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 112, 192)
    line.line.fill.background()

    # Content
    if content_items:
        if image_path and image_pos:
            tw = text_width if text_width else 4.5
            content_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(tw), Inches(5.5))
        else:
            content_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(9.4), Inches(5.5))

        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Check formatting type
            if item.startswith("  -"):
                # Sub-bullet (indented)
                text = "    - " + item.strip()[2:].strip()
                p.level = 1
                add_formatted_text(p, text, font_size - 2)
            elif item.startswith(">>"):
                # Highlighted/emphasized item
                text = "→ " + item[2:].strip()
                add_formatted_text(p, text, font_size, is_highlight=True)
            else:
                # Regular bullet with possible inline bold
                text = "• " + item if not item.startswith("•") else item
                add_formatted_text(p, text, font_size)

            p.space_after = Pt(6)

    # Image
    if image_path and os.path.exists(image_path):
        if image_pos:
            left, top = image_pos[0], image_pos[1]
            width = image_pos[2] if len(image_pos) > 2 and image_pos[2] else None
            height = image_pos[3] if len(image_pos) > 3 and image_pos[3] else None
            slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                    width=Inches(width) if width else None,
                                    height=Inches(height) if height else None)

    return slide


def add_image_slide(prs, title, image_path, caption="", image_size=None):
    """Add a slide with primarily an image."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Image
    if os.path.exists(image_path):
        if image_size:
            left, top, width, height = image_size
            slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                    width=Inches(width) if width else None,
                                    height=Inches(height) if height else None)
        else:
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), width=Inches(9))

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.8), Inches(9.4), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_image_slide(prs, title, img1_path, img2_path, label1="", label2=""):
    """Add a slide with two images side by side."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Left image
    if os.path.exists(img1_path):
        slide.shapes.add_picture(img1_path, Inches(0.3), Inches(1.2), width=Inches(4.5))

    # Right image
    if os.path.exists(img2_path):
        slide.shapes.add_picture(img2_path, Inches(5.2), Inches(1.2), width=Inches(4.5))

    # Labels
    if label1:
        lbl1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.8), Inches(4.5), Inches(0.4))
        tf = lbl1.text_frame
        p = tf.paragraphs[0]
        p.text = label1
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER

    if label2:
        lbl2 = slide.shapes.add_textbox(Inches(5.2), Inches(5.8), Inches(4.5), Inches(0.4))
        tf = lbl2.text_frame
        p = tf.paragraphs[0]
        p.text = label2
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_table_slide(prs, title, headers, rows, subtitle=""):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Subtitle
    top_offset = 1.2
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(9.4), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        top_offset = 1.6

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(top_offset),
                                   Inches(9), Inches(0.5 * num_rows)).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    return slide


def add_equation_slide(prs, title, equation_text, explanation_items=None):
    """Add a slide with equation and explanation."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Equation box
    eq_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(1.2))
    tf = eq_box.text_frame
    p = tf.paragraphs[0]
    p.text = equation_text
    p.font.size = Pt(22)
    p.alignment = PP_ALIGN.CENTER

    # Explanation
    if explanation_items:
        exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(4))
        tf = exp_box.text_frame
        for i, item in enumerate(explanation_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if item.startswith("**") and item.endswith("**"):
                p.text = item[2:-2]
                p.font.bold = True
                p.font.size = Pt(16)
            else:
                p.text = "• " + item
                p.font.size = Pt(16)
            p.space_after = Pt(8)

    return slide


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # Slide 1: Title
    # =========================================================================
    add_title_slide(prs,
                   "HyperCOT-MS",
                   "Hypergraph Co-Optimal Transport\nfor Morse-Smale Complex Comparison")

    # =========================================================================
    # Slide 2: Problem Statement
    # =========================================================================
    add_content_slide(prs, "Problem: Comparing Morse-Smale Complexes", [
        "**Goal**: Measure similarity between two MS complexes",
        "**Application**: Compare clean vs. noisy scalar fields",
        "**Existing methods**:",
        "  - Wasserstein Distance (WD): Uses scalar field values only",
        "  - Gromov-Wasserstein (GWD): Uses graph structure only",
        "**Limitation**: Neither captures the full MS complex structure",
        "  - MS complex has regions (not just points and edges)",
        "  - Regions form a hypergraph structure"
    ])

    # =========================================================================
    # Slide 3: Input Data
    # =========================================================================
    add_two_image_slide(prs, "Input: Clean vs. Noisy MS Complex",
                       os.path.join(BASE_PATH, "clean_input.png"),
                       os.path.join(BASE_PATH, "noise_input.png"),
                       "Clean: 49 CPs, 36 regions",
                       "Noisy: 65 CPs, 47 regions")

    # =========================================================================
    # Slide 4: Proposed Approach
    # =========================================================================
    add_content_slide(prs, "Proposed Approach: HyperCOT", [
        "**Key Idea**: Model MS complex as a hypergraph",
        "  - Nodes = Critical points",
        "  - Hyperedges = Regions (bounded by multiple CPs)",
        "**HyperCOT**: Hypergraph Co-Optimal Transport",
        "  - Jointly optimize node coupling (π) and hyperedge coupling (ξ)",
        "  - Captures both point-level and region-level correspondence"
    ])

    # =========================================================================
    # Slide 5: Challenges / Outline
    # =========================================================================
    add_content_slide(prs, "Key Challenges", [
        "**Challenge 1: Hypergraph Construction**",
        "  - How to extract hyperedges from MS complex?",
        "  - What defines the boundary of a region?",
        "**Challenge 2: Distance from CP to Hyperedge**",
        "  - Hyperedge is an area, not a point",
        "  - Need to define hypernetwork function ω(v, e)",
        "**Challenge 3: Measure Definition**",
        "  - Node measure μ: probability over CPs",
        "  - Hyperedge measure ν: probability over regions",
        ">>This talk focuses on these three implementation challenges"
    ])

    # =========================================================================
    # Slide 6: HyperCOT Formulation (Overview)
    # =========================================================================
    add_equation_slide(prs, "HyperCOT: Objective Function",
                      "min     Σ |ω₁(v,e) - ω₂(v',e')|² · π(v,v') · ξ(e,e')\n  π,ξ",
                      [
                          "π ∈ Π(μ₁, μ₂): Node coupling with marginals μ₁, μ₂",
                          "ξ ∈ Π(ν₁, ν₂): Hyperedge coupling with marginals ν₁, ν₂",
                          "ω(v, e): Hypernetwork function (distance from CP v to region e)",
                          "**Three required components:**",
                          "μ (node measure), ν (hyperedge measure), ω (structure)"
                      ])

    # =========================================================================
    # Slide 7: HyperCOT Formal Definition
    # =========================================================================
    add_content_slide(prs, "HyperCOT: Formal Definition", [
        "**Hypergraph representation**:",
        "  - H = (V, E, ω, μ, ν) where V = CPs, E = regions",
        "**Objective function**:",
        "  - COOT(H₁,H₂) = min_{π,ξ} Σ_{v,v'} Σ_{e,e'} |ω₁(v,e) - ω₂(v',e')|² π(v,v') ξ(e,e')",
        "**Constraints**:",
        "  - π ∈ Π(μ₁, μ₂):  π𝟙 = μ₁,  πᵀ𝟙 = μ₂",
        "  - ξ ∈ Π(ν₁, ν₂):  ξ𝟙 = ν₁,  ξᵀ𝟙 = ν₂",
        "**Notation**:",
        "  - π ∈ ℝ^{|V₁|×|V₂|}: Node (CP) coupling matrix",
        "  - ξ ∈ ℝ^{|E₁|×|E₂|}: Hyperedge (region) coupling matrix",
        "  - ω(v, e): Distance from CP v to region e"
    ], font_size=17)

    # =========================================================================
    # Slide 8: Challenge 1 - Hypergraph Construction
    # =========================================================================
    add_content_slide(prs, "Challenge 1: Hypergraph Construction", [
        "**Goal**: Convert MS complex regions into hyperedges",
        "**Observation**: Each MS region is bounded by exactly:",
        "  - 1 minimum",
        "  - 2 saddles",
        "  - 1 maximum",
        "**Hyperedge definition**:",
        "  - hyperedge(e) = {min, saddle₁, saddle₂, max}",
        "**Input**: TTK segmentation data",
        "**Output**: hypergraph_clean.csv, hypergraph_noisy.csv"
    ], os.path.join(BASE_PATH, "coot_input.png"), (5.2, 1.5, 4.3, None))

    # =========================================================================
    # Slide 8: Hypergraph Construction - Implementation
    # =========================================================================
    add_content_slide(prs, "Hypergraph Construction: Implementation", [
        "**Step 1: Parse TTK segmentation**",
        "  - Each mesh vertex → (DescendingManifold, AscendingManifold)",
        "  - DescendingManifold = minimum ID",
        "  - AscendingManifold = maximum ID (with offset)",
        "**Step 2: Group by region**",
        "  - Unique (min_id, max_id) pairs define regions",
        "**Step 3: Find boundary saddles**",
        "  - Identify saddles on region boundaries",
        "**Step 4: Build hyperedge**",
        "  - hyperedge = [min_id, saddle₁, saddle₂, max_id]"
    ])

    # =========================================================================
    # Slide 9: Challenge 2 - Distance CP to Hyperedge
    # =========================================================================
    add_content_slide(prs, "Challenge 2: Distance from CP to Hyperedge", [
        "**Problem**: ω(v, e) requires distance from point to region",
        "  - Region is an area, not a single point",
        "  - Cannot directly compute point-to-area distance",
        "**Solution**: Virtual Center (VC) representation",
        "  - Approximate each region with a representative point",
        "  - ω(v, e) = shortest path from CP v to VC of region e",
        "**Implementation**:",
        "  - Step 1: Generate virtual centers",
        "  - Step 2: Build augmented graph",
        "  - Step 3: Compute shortest paths via Dijkstra"
    ])

    # =========================================================================
    # Slide 10: Virtual Center Generation
    # =========================================================================
    add_content_slide(prs, "Virtual Center Generation", [
        "**Goal**: Find a representative point for each region",
        "**Method**: Line intersection",
        "  - Line 1: minimum → maximum",
        "  - Line 2: saddle₁ → saddle₂",
        "  - VC = intersection point",
        "**Computation**:",
        "  - Solve in 2D (x, y coordinates)",
        "  - Interpolate z from boundary CPs",
        "**Output**: One VC per region"
    ], os.path.join(BASE_PATH, "clean_vc_adjacency.png"), (5, 2.0, 4.8, None), text_width=4.8)

    # =========================================================================
    # Slide 11: Augmented Graph Construction
    # =========================================================================
    add_content_slide(prs, "Augmented Graph Construction", [
        "**Nodes**: Critical Points ∪ Virtual Centers",
        "**Edge Type 1: CP ↔ CP**",
        "  - From separatrices (TTK output)",
        "  - Weight: Euclidean distance",
        "**Edge Type 2: CP ↔ VC**",
        "  - Connect boundary CPs to their region's VC",
        "  - Weight: Euclidean distance",
        "**Edge Type 3: VC ↔ VC**",
        "  - Connect adjacent regions (share boundary CPs)",
        "  - Weight: Euclidean distance between VCs"
    ])

    # =========================================================================
    # Slide 12: VC-VC Adjacency
    # =========================================================================
    add_content_slide(prs, "VC-VC Adjacency", [
        "**Definition**: Two regions are adjacent if they share boundary CPs",
        "**Strong adjacency**: Regions sharing 2+ CPs",
        "**Statistics**:",
        "  - Clean: 36 regions, 60 strong adjacencies",
        "  - Noisy: 47 regions, 74 strong adjacencies",
        "**Visualization**: Orange edges in augmented graph"
    ], os.path.join(BASE_PATH, "clean_vc_adjacency.png"), (4.8, 1.8, 5, None), text_width=4.5)

    # =========================================================================
    # Slide 13: Shortest Path Computation
    # =========================================================================
    add_content_slide(prs, "Hypernetwork Function ω: Shortest Path", [
        "**Definition**: ω(v, e) = shortest path from CP v to VC of region e",
        "**Algorithm**: Dijkstra on augmented graph",
        "**Path can traverse**:",
        "  - CP → CP edges (along separatrices)",
        "  - CP → VC edges (enter a region)",
        "  - VC → VC edges (cross between regions)",
        "**Output**: ω matrix of size (n_CPs × n_regions)",
        "**Example**: ω(45, 20) = 66.98"
    ])

    # =========================================================================
    # Slide 14: Shortest Path Example
    # =========================================================================
    add_image_slide(prs, "Shortest Path Example",
                   os.path.join(BASE_PATH, "clean_vc_adjacency.png"),
                   "Left: VC generation. Middle: Augmented graph. Right: Path CP 45 → VC28 → VC20 (ω = 66.98)",
                   (0.3, 1.0, 9.4, None))

    # =========================================================================
    # Slide 15: Challenge 3 - Measure Definition
    # =========================================================================
    add_content_slide(prs, "Challenge 3: Measure Definition", [
        "**Node measure μ (existing work)**:",
        "  - Extended persistence on CP graph",
        "  - Persistence image representation",
        "  - μ(v) = normalized contribution of CP v",
        "**Hyperedge measure ν (this work)**:",
        "  - Based on boundary CP measures",
        "  - ν(e) = Σ_{v ∈ boundary(e)} μ(v) / Z",
        "  - Intuition: Region importance = sum of boundary CP importance"
    ])

    # =========================================================================
    # Slide 16: Node Measure Visualization
    # =========================================================================
    add_image_slide(prs, "Node Measure μ: Extended Persistence",
                   os.path.join(BASE_PATH, "clean_extended.png"),
                   "Left: Extended persistence diagram. Middle: Persistence image. Right: Node measure μ distribution.",
                   (0.3, 1.0, 9.4, None))

    # =========================================================================
    # Slide 17: HyperCOT Optimization
    # =========================================================================
    add_content_slide(prs, "HyperCOT Optimization", [
        "**Algorithm**: Alternating Sinkhorn iteration",
        "**Initialization**:",
        "  - π = μ₁ ⊗ μ₂ (outer product)",
        "  - ξ = ν₁ ⊗ ν₂ (outer product)",
        "**Iteration**:",
        "  - Fix ξ, update π via Sinkhorn",
        "  - Fix π, update ξ via Sinkhorn",
        "**Entropy regularization**: ε = 0.05",
        "**Output**:",
        "  - π: 49 × 65 node coupling matrix",
        "  - ξ: 36 × 47 hyperedge coupling matrix"
    ])

    # =========================================================================
    # Slide 18: Results - Comparison
    # =========================================================================
    add_table_slide(prs, "Results: Method Comparison",
                   ["Method", "Distance", "Type Preservation", "Structure Used"],
                   [
                       ["WD", "0.0402", "90.0%", "Scalar values only"],
                       ["GWD", "0.0037", "72.5%", "Graph structure only"],
                       ["HyperCOT", "0.0024", "85.7%", "Hypergraph structure"]
                   ],
                   "Lower distance = better matching; Higher type preservation = CPs matched to same type")

    # =========================================================================
    # Slide 19: Results - HyperCOT Visualization
    # =========================================================================
    add_image_slide(prs, "Results: HyperCOT Correspondence",
                   os.path.join(BASE_PATH, "hypercot_detailed_correspondence.png"),
                   "Left: Spatial correspondence. Middle: Node coupling π. Right: Hyperedge coupling ξ.",
                   (0.3, 1.0, 9.4, None))

    # =========================================================================
    # Slide 20: Discussion
    # =========================================================================
    add_content_slide(prs, "Discussion", [
        "**Advantages of HyperCOT**:",
        "  - Captures region-level structure via hypergraph",
        "  - Joint optimization of node and hyperedge coupling",
        "  - Virtual centers enable point-to-region distance",
        "**Limitations**:",
        "  - VC may not lie inside the region (approximation)",
        "  - Computational cost of alternating optimization",
        "**Future Work**:",
        "  - Apply to real-world datasets (combustion, climate)",
        "  - Improve VC computation (centroid-based)",
        "  - Extend to 3D MS complexes"
    ])

    # =========================================================================
    # Slide 21: Summary
    # =========================================================================
    add_content_slide(prs, "Summary", [
        "**HyperCOT-MS**: Compare MS complexes using hypergraph OT",
        "**Key contributions**:",
        "  - Hypergraph construction from TTK segmentation",
        "  - Virtual center for region representation",
        "  - Augmented graph with VC-VC adjacency",
        "  - Shortest path for hypernetwork function ω",
        "**Results**:",
        "  - Lowest distance among WD, GWD, HyperCOT",
        "  - Good type preservation (85.7%)",
        "**Code**: github.com/gm3g11/HyperCOT-MS"
    ])

    # Save
    output_path = os.path.join(BASE_PATH, "HyperCOT-MS_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    create_presentation()
